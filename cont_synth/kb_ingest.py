"""Knowledge Base ingestion pipeline.

Handles text extraction, chunking, embedding, and storage for the RAG system.
Also provides the retrieval function used during synthesis.
"""
import io
import re
from typing import Optional

import PyPDF2
import reflex as rx
from sqlmodel import select

from .models import WorkspaceDocument, DocumentChunk

# ---------------------------------------------------------------------------
# Text Sanitization
# ---------------------------------------------------------------------------

# Matches sequences of 3+ single letters each separated by exactly one space.
# e.g. "P r o d u c t" (7 letters) or "A B C" (3 letters).
# Requires ≥3 letters ({2,} means the "space+letter" group repeats ≥2 times)
# to avoid false positives on natural two-letter combos like "a b" or "I m".
_SPACED_LETTERS_RE = re.compile(r'(?<![A-Za-z])([A-Za-z])( [A-Za-z]){2,}(?![A-Za-z])')


def sanitize_text(text: str) -> str:
    """Fix PDF kerning artifacts where letter-spacing is parsed as literal spaces.

    PDF extractors often interpret visual letter-spacing in headings as actual
    space characters, producing output like "P r o d u c t  C o m p o n e n t s"
    instead of "Product Components".

    Strategy:
      1. Match any run of 3+ single letters each separated by exactly one space.
      2. Remove the spaces within that run, condensing back into a word.
      3. Normalize leftover horizontal whitespace runs (the double-space word
         boundaries between formerly-spaced words) down to a single space.
         Newlines are intentionally preserved so paragraph structure survives.
    """
    # Step 1: collapse spaced-letter runs into solid words
    text = _SPACED_LETTERS_RE.sub(lambda m: m.group(0).replace(' ', ''), text)
    # Step 2: collapse runs of spaces/tabs (but NOT newlines) left behind
    text = re.sub(r'[ \t]{2,}', ' ', text)
    return text

# Lazy-loaded SentenceTransformer model (avoid slow startup)
_st_model = None


def _get_st_model():
    """Load and cache the SentenceTransformer model on first call.

    Passes HF_TOKEN from the environment if present, which suppresses the
    Hugging Face anonymous-access warning. Silently ignored when absent.
    """
    global _st_model
    if _st_model is None:
        import os
        from sentence_transformers import SentenceTransformer
        hf_token = os.getenv("HF_TOKEN") or None
        _st_model = SentenceTransformer('all-MiniLM-L6-v2', token=hf_token)
    return _st_model


# ---------------------------------------------------------------------------
# Text Extraction
# ---------------------------------------------------------------------------

def extract_text(file_data: bytes, filename: str) -> str:
    """Extract raw text from a PDF, TXT, DOCX, or PPTX file."""
    lower = filename.lower()

    if lower.endswith('.pdf'):
        reader = PyPDF2.PdfReader(io.BytesIO(file_data))
        pages = [page.extract_text() for page in reader.pages if page.extract_text()]
        return "\n".join(pages)

    if lower.endswith('.docx'):
        from docx import Document
        doc = Document(io.BytesIO(file_data))
        return "\n".join(para.text for para in doc.paragraphs if para.text)

    if lower.endswith('.pptx'):
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)

    if lower.endswith('.doc'):
        raise ValueError(
            "Old .doc format is not supported. Please open the file in Word and save it as .docx, then re-upload."
        )

    if lower.endswith('.ppt'):
        raise ValueError(
            "Old .ppt format is not supported. Please open the file in PowerPoint and save it as .pptx, then re-upload."
        )

    # Fallback: treat as UTF-8 text
    return file_data.decode('utf-8', errors='replace')


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> list[str]:
    """Split text into overlapping character-level chunks.

    Args:
        text: The full document text.
        chunk_size: Target characters per chunk.
        overlap: Characters to repeat at the start of each subsequent chunk.

    Returns:
        List of non-empty text chunks.
    """
    text = text.strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= len(text):
            break
        start = end - overlap
    return chunks


# ---------------------------------------------------------------------------
# Embedding
# ---------------------------------------------------------------------------

def embed_chunks(chunks: list[str]) -> list[list[float]]:
    """Vectorise a list of text chunks using all-MiniLM-L6-v2 (384 dims)."""
    model = _get_st_model()
    embeddings = model.encode(chunks, show_progress_bar=False)
    return [emb.tolist() for emb in embeddings]


# ---------------------------------------------------------------------------
# Full Ingestion
# ---------------------------------------------------------------------------

def ingest_document(file_data: bytes, filename: str, product_id: int) -> WorkspaceDocument:
    """Extract, chunk, embed, and persist a document to the database.

    Returns the created WorkspaceDocument record.
    """
    text = sanitize_text(extract_text(file_data, filename))
    chunks = chunk_text(text)

    if not chunks:
        raise ValueError(f"No text could be extracted from '{filename}'.")

    embeddings = embed_chunks(chunks)

    with rx.session() as session:
        doc = WorkspaceDocument(product_id=product_id, filename=filename, file_data=file_data)
        session.add(doc)
        session.commit()
        session.refresh(doc)

        for chunk_text_str, embedding in zip(chunks, embeddings):
            chunk = DocumentChunk(
                document_id=doc.id,
                chunk_text=chunk_text_str,
                embedding=embedding,
            )
            session.add(chunk)
        session.commit()
        session.refresh(doc)

    return doc


# ---------------------------------------------------------------------------
# RAG Retrieval (used by the Synthesis engine)
# ---------------------------------------------------------------------------

def get_rag_context(transcript: str, product_id: int, top_k: int = 5) -> str:
    """Build a workspace context string for the synthesis prompt via RAG.

    Steps:
      1. Summarise the transcript with gemini-flash (focuses the query vector).
      2. Embed the summary with all-MiniLM-L6-v2.
      3. Cosine-similarity search DocumentChunk for the top-k matches.
      4. Return the chunks concatenated with separators.

    Returns an empty string if there are no documents in the KB or any step fails.
    """
    # 1. Check whether there are any documents for this product first (fast path)
    with rx.session() as session:
        doc_ids = [
            row.id for row in session.exec(
                select(WorkspaceDocument).where(WorkspaceDocument.product_id == product_id)
            ).all()
        ]
    if not doc_ids:
        return ""

    # 2. Summarise the transcript to create a focused query
    from .state.core import flash_model, load_prompt

    summary_prompt = load_prompt("rag_summarize.txt").format(transcript=transcript[:4000])
    summary_resp = flash_model.generate_content(summary_prompt)
    summary = summary_resp.text.strip()

    # 3. Embed the summary
    model = _get_st_model()
    query_vector = model.encode(summary).tolist()

    # 4. Cosine similarity search via pgvector
    with rx.session() as session:
        # pgvector's <=> operator computes cosine distance; order ascending = most similar first
        from sqlalchemy import text as sa_text
        rows = session.exec(
            sa_text(
                "SELECT chunk_text FROM documentchunk "
                "WHERE document_id = ANY(:doc_ids) "
                "ORDER BY embedding <=> CAST(:qv AS vector) "
                "LIMIT :k"
            ).bindparams(
                doc_ids=doc_ids,
                qv=str(query_vector),
                k=top_k,
            )
        ).all()

    if not rows:
        return ""

    chunks = [row[0] for row in rows]
    return "\n\n---\n\n".join(chunks)


# ---------------------------------------------------------------------------
# Prep-page RAG Retrieval (query by persona + opportunities, no transcript)
# ---------------------------------------------------------------------------

def get_prep_rag_context(query_text: str, product_id: int, top_k: int = 5) -> str:
    """Retrieve workspace context relevant to a prep query (no transcript available).

    Unlike get_rag_context, this skips the LLM summarisation step and embeds
    the caller-supplied query_text directly.  The caller should combine the
    target persona name and the selected opportunity statements so the vector
    search targets the right product area.

    Returns an empty string if no documents exist or any step fails.
    """
    with rx.session() as session:
        doc_ids = [
            row.id for row in session.exec(
                select(WorkspaceDocument).where(WorkspaceDocument.product_id == product_id)
            ).all()
        ]
    if not doc_ids:
        return ""

    model = _get_st_model()
    query_vector = model.encode(query_text).tolist()

    with rx.session() as session:
        from sqlalchemy import text as sa_text
        rows = session.exec(
            sa_text(
                "SELECT chunk_text FROM documentchunk "
                "WHERE document_id = ANY(:doc_ids) "
                "ORDER BY embedding <=> CAST(:qv AS vector) "
                "LIMIT :k"
            ).bindparams(
                doc_ids=doc_ids,
                qv=str(query_vector),
                k=top_k,
            )
        ).all()

    if not rows:
        return ""

    chunks = [row[0] for row in rows]
    return "\n\n---\n\n".join(chunks)


# ---------------------------------------------------------------------------
# Re-processing (sanitize + re-embed existing stored chunks)
# ---------------------------------------------------------------------------

def reprocess_document_chunks(doc_id: int) -> int:
    """Sanitize chunk_text and re-embed all chunks for a given WorkspaceDocument.

    Because we don't store original file bytes, we sanitize the already-stored
    chunk_text in place and regenerate embeddings from the cleaned text.

    Returns the number of chunks updated.
    """
    with rx.session() as session:
        chunks = session.exec(
            select(DocumentChunk).where(DocumentChunk.document_id == doc_id)
        ).all()
        if not chunks:
            return 0

        clean_texts = [sanitize_text(c.chunk_text) for c in chunks]
        new_embeddings = embed_chunks(clean_texts)

        for chunk, clean_text, embedding in zip(chunks, clean_texts, new_embeddings):
            chunk.chunk_text = clean_text
            chunk.embedding = embedding
            session.add(chunk)
        session.commit()

    return len(chunks)
